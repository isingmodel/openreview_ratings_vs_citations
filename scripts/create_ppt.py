"""Generate a one-page PowerPoint presentation for the research project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


def create_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Colors
    DARK_BG = RGBColor(26, 26, 46)  # #1a1a2e
    ACCENT = RGBColor(0, 212, 255)  # #00d4ff
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(255, 215, 0)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = "📊 Do Peer Review Scores Predict Scientific Impact?"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.paragraphs[0].text = "Empirical Analysis of ICLR Papers (2017-2020)  •  PyCon Korea 2022"
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT
    
    # Key Finding Box
    finding_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(5), Inches(1.2)
    )
    finding_box.fill.solid()
    finding_box.fill.fore_color.rgb = RGBColor(45, 45, 68)
    finding_box.line.fill.background()
    
    finding_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.85), Inches(4.6), Inches(1))
    finding_frame = finding_text.text_frame
    finding_frame.word_wrap = True
    p = finding_frame.paragraphs[0]
    p.text = "🔍 Key Finding"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    p2 = finding_frame.add_paragraph()
    p2.text = "Correlation between ratings and citations is declining over time (r: 0.40 → 0.13)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    
    # Results Table
    table_data = [
        ["Year", "Papers", "r", "p-value"],
        ["2017", "245", "0.40", "6.9e-11"],
        ["2018", "425", "0.37", "1.6e-15"],
        ["2019", "502", "0.19", "1.2e-05"],
        ["2020", "687", "0.13", "8.9e-04"],
    ]
    
    table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(3.1), Inches(5), Inches(1.8)).table
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(60, 60, 90)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.fore_color.rgb = RGBColor(40, 40, 60)
    
    # Add images (2x2 grid on right side)
    img_width = Inches(3.2)
    img_height = Inches(2.4)
    
    figs_dir = Path(__file__).parent.parent / "figs"
    
    # Top row
    slide.shapes.add_picture(str(figs_dir / "Log_Citation_vs_Review_Rating_ICLR_2017.png"),
                             Inches(6), Inches(1.5), img_width, img_height)
    slide.shapes.add_picture(str(figs_dir / "Log_Citation_vs_Review_Rating_ICLR_2018.png"),
                             Inches(9.5), Inches(1.5), img_width, img_height)
    # Bottom row
    slide.shapes.add_picture(str(figs_dir / "Log_Citation_vs_Review_Rating_ICLR_2019.png"),
                             Inches(6), Inches(4.1), img_width, img_height)
    slide.shapes.add_picture(str(figs_dir / "Log_Citation_vs_Review_Rating_ICLR_2020.png"),
                             Inches(9.5), Inches(4.1), img_width, img_height)
    
    # Methodology note
    method_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5), Inches(1))
    method_frame = method_box.text_frame
    method_frame.word_wrap = True
    p = method_frame.paragraphs[0]
    p.text = "📐 Methodology: Pearson correlation on log(citations+1)"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(180, 180, 180)
    
    p2 = method_frame.add_paragraph()
    p2.text = "Log transform reduces outlier dominance in heavy-tailed citation distributions"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(140, 140, 140)
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_frame.paragraphs[0].text = "github.com/isingmodel/openreview_ratings_vs_citations"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    # Save
    output_path = Path(__file__).parent.parent / "OpenReview_Research_Intro.pptx"
    prs.save(str(output_path))
    print(f"Saved presentation to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
